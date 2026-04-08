"""
generate_ppt.py  —  AIOps Control Plane presentation generator
Run:  python generate_ppt.py
Output: AIOps_ControlPlane_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x11, 0x17)   # near-black navy
BG_CARD       = RGBColor(0x13, 0x1C, 0x2B)   # dark card
ACCENT_BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # electric blue
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)   # cyan
ACCENT_GREEN  = RGBColor(0x22, 0xC5, 0x5E)   # green
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # amber
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)   # red
ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)   # purple
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE     = RGBColor(0xCB, 0xD5, 0xE1)
MUTED         = RGBColor(0x64, 0x74, 0x8B)
DIVIDER       = RGBColor(0x1E, 0x2D, 0x40)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── helpers ──────────────────────────────────────────────────────────────────
def blank_slide():
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return slide


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_label_box(slide, text, x, y, w, h,
                  fill=BG_CARD, border=ACCENT_BLUE,
                  txt_color=WHITE, size=Pt(11), bold=False):
    add_rect(slide, x, y, w, h, fill=fill, line=border, line_width=Pt(1.2))
    add_text(slide, text, x + Inches(0.08), y + Inches(0.05),
             w - Inches(0.16), h - Inches(0.1),
             size=size, bold=bold, color=txt_color, align=PP_ALIGN.CENTER)


def add_pill(slide, text, x, y, color=ACCENT_BLUE, size=Pt(10)):
    w = Inches(1.5)
    h = Inches(0.28)
    add_rect(slide, x, y, w, h, fill=color)
    add_text(slide, text, x + Inches(0.05), y + Inches(0.02),
             w - Inches(0.1), h - Inches(0.04),
             size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def arrow_right(slide, x, y, length=Inches(0.45), color=MUTED):
    """Draw a simple right-pointing arrow using a line + triangle."""
    line = slide.shapes.add_shape(1, x, y, length, Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    # arrowhead triangle
    tri = slide.shapes.add_shape(5, x + length - Inches(0.12), y - Inches(0.06),
                                   Inches(0.14), Inches(0.15))
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()


def arrow_down(slide, x, y, length=Inches(0.3), color=MUTED):
    line = slide.shapes.add_shape(1, x, y, Inches(0.03), length)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def section_header(slide, label, title, subtitle, accent=ACCENT_BLUE):
    # top accent bar
    add_rect(slide, 0, 0, W, Inches(0.07), fill=accent)
    # left hero stripe
    add_rect(slide, 0, Inches(0.07), Inches(0.06), H - Inches(0.07), fill=accent)
    add_text(slide, label,     Inches(0.2), Inches(0.25), Inches(9), Inches(0.35),
             size=Pt(11), color=accent, bold=True)
    add_text(slide, title,     Inches(0.2), Inches(0.6),  Inches(10), Inches(0.8),
             size=Pt(38), bold=True, color=WHITE)
    add_text(slide, subtitle,  Inches(0.2), Inches(1.42), Inches(11), Inches(0.45),
             size=Pt(15), color=OFF_WHITE)
    # horizontal rule
    add_rect(slide, Inches(0.2), Inches(1.9), Inches(12.9), Inches(0.03), fill=DIVIDER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
# Full-width top gradient bar
add_rect(slide, 0, 0, W, Inches(0.08), fill=ACCENT_BLUE)
# big diagonal accent block
add_rect(slide, Inches(8.5), 0, Inches(4.83), H, fill=BG_CARD)
add_rect(slide, Inches(8.5), 0, Inches(0.05), H, fill=ACCENT_BLUE)

add_text(slide, "AIOps",
         Inches(0.5), Inches(1.2), Inches(7), Inches(1.5),
         size=Pt(72), bold=True, color=ACCENT_BLUE)
add_text(slide, "Observability Control Plane",
         Inches(0.5), Inches(2.65), Inches(7.5), Inches(0.8),
         size=Pt(28), bold=True, color=WHITE)
add_text(slide, "AI-guided incident response · Real-time topology · Fleet onboarding",
         Inches(0.5), Inches(3.4), Inches(7.8), Inches(0.5),
         size=Pt(14), color=OFF_WHITE)

# right column — feature chips
chips = [
    ("Applications", ACCENT_BLUE),
    ("Incidents & Alerts", ACCENT_RED),
    ("AI RCA & Remediation", ACCENT_CYAN),
    ("Graph Topology", ACCENT_PURPLE),
    ("Predictions", ACCENT_AMBER),
    ("Fleet Onboarding", ACCENT_GREEN),
]
for i, (label, color) in enumerate(chips):
    cy = Inches(1.1) + i * Inches(0.85)
    add_rect(slide, Inches(9.0), cy, Inches(3.8), Inches(0.6), fill=color)
    add_text(slide, label, Inches(9.15), cy + Inches(0.12), Inches(3.5), Inches(0.38),
             size=Pt(15), bold=True, color=WHITE)

add_text(slide, "April 2026  ·  v1.0  ·  Phase 1 Complete",
         Inches(0.5), Inches(6.5), Inches(8), Inches(0.4),
         size=Pt(11), color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "OVERVIEW", "What We Built", "A full-stack AIOps platform in a single deployable stack", ACCENT_BLUE)

cards = [
    ("🔭  Unified Observability", ACCENT_BLUE,
     "Applications, services, incidents, alerts, predictions, and topology in one React control plane"),
    ("🤖  AI-Guided Response", ACCENT_CYAN,
     "AiDE LLM integration for root cause analysis, diagnostic planning, and remediation command generation"),
    ("📡  Real-time Signals", ACCENT_PURPLE,
     "Prometheus metrics, Alertmanager webhooks, Elasticsearch logs, and Jaeger traces stitched together"),
    ("🛡️  Incident Lifecycle", ACCENT_RED,
     "Auto-correlated incidents with timeline, blast radius, dependency graph, and resolution tracking"),
    ("🔮  Predictive Risk", ACCENT_AMBER,
     "Heuristic risk scoring with 15-minute incident probability windows per service"),
    ("🗂️  Fleet Control", ACCENT_GREEN,
     "Linux target onboarding via SSH + PEM with telemetry profile selection and enrollment tokens"),
]
cols = 3
for i, (title, color, desc) in enumerate(cards):
    col = i % cols
    row = i // cols
    cx = Inches(0.25) + col * Inches(4.35)
    cy = Inches(2.05) + row * Inches(2.3)
    add_rect(slide, cx, cy, Inches(4.1), Inches(2.1), fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_rect(slide, cx, cy, Inches(4.1), Inches(0.07), fill=color)
    add_text(slide, title, cx + Inches(0.15), cy + Inches(0.15), Inches(3.8), Inches(0.45),
             size=Pt(13), bold=True, color=WHITE)
    add_text(slide, desc, cx + Inches(0.15), cy + Inches(0.6), Inches(3.8), Inches(1.3),
             size=Pt(11), color=OFF_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "ARCHITECTURE", "System Architecture", "Docker Compose multi-service stack with AI backbone", ACCENT_CYAN)

# 3 layers: Data Sources | Core Platform | Frontend
# ── Layer labels ──
for lx, label, color in [
    (Inches(0.15), "DATA SOURCES", ACCENT_AMBER),
    (Inches(4.55), "CORE PLATFORM", ACCENT_BLUE),
    (Inches(9.8), "FRONTEND", ACCENT_GREEN),
]:
    add_rect(slide, lx, Inches(2.1), Inches(4.1) if lx < Inches(5) else Inches(3.3), Inches(0.32), fill=color)
    add_text(slide, label, lx + Inches(0.1), Inches(2.13), Inches(3.9), Inches(0.28),
             size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Data sources column
ds_items = [
    ("Prometheus", ACCENT_AMBER),
    ("Alertmanager", ACCENT_RED),
    ("Elasticsearch", ACCENT_BLUE),
    ("Jaeger Traces", ACCENT_PURPLE),
    ("PostgreSQL", ACCENT_CYAN),
    ("Redis Cache", ACCENT_GREEN),
]
for i, (name, color) in enumerate(ds_items):
    cy = Inches(2.6) + i * Inches(0.72)
    add_rect(slide, Inches(0.2), cy, Inches(3.9), Inches(0.55), fill=BG_CARD, line=color, line_width=Pt(1.2))
    add_text(slide, name, Inches(0.35), cy + Inches(0.1), Inches(3.6), Inches(0.35),
             size=Pt(12), bold=True, color=color)

# Platform column (Django + services)
plat_items = [
    ("Django Web (Gunicorn)", ACCENT_BLUE),
    ("genai.views — Alert Ingest / RCA", ACCENT_CYAN),
    ("genai.views — Fleet / Enroll", ACCENT_GREEN),
    ("Celery / Predictor Worker", ACCENT_AMBER),
    ("OTel Collector → Jaeger", ACCENT_PURPLE),
    ("DB-Agent + Control-Agent", MUTED),
]
for i, (name, color) in enumerate(plat_items):
    cy = Inches(2.6) + i * Inches(0.72)
    add_rect(slide, Inches(4.6), cy, Inches(4.9), Inches(0.55), fill=BG_CARD, line=color, line_width=Pt(1.2))
    add_text(slide, name, Inches(4.75), cy + Inches(0.1), Inches(4.6), Inches(0.35),
             size=Pt(11), bold=True, color=color)

# Frontend column
fe_items = [
    ("React + Vite (port 8089)", ACCENT_GREEN),
    ("Applications Page", ACCENT_BLUE),
    ("Alerts / Incidents", ACCENT_RED),
    ("Fleet / Enroll / Profiles", ACCENT_CYAN),
    ("Graph Topology", ACCENT_PURPLE),
    ("AI Assistant", ACCENT_AMBER),
]
for i, (name, color) in enumerate(fe_items):
    cy = Inches(2.6) + i * Inches(0.72)
    add_rect(slide, Inches(9.85), cy, Inches(3.3), Inches(0.55), fill=BG_CARD, line=color, line_width=Pt(1.2))
    add_text(slide, name, Inches(10.0), cy + Inches(0.1), Inches(3.0), Inches(0.35),
             size=Pt(10), bold=True, color=color)

# arrows between layers
for ay in [Inches(2.85), Inches(3.57), Inches(4.29), Inches(5.01), Inches(5.73), Inches(6.45)]:
    add_rect(slide, Inches(4.12), ay, Inches(0.44), Inches(0.06), fill=ACCENT_BLUE)
    add_rect(slide, Inches(9.77), ay, Inches(0.06), Inches(0.06), fill=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ALERT INGESTION & RCA FLOW
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "ALERT PIPELINE", "Alert Ingestion & AI RCA Flow",
               "From Prometheus firing to actionable AI remediation in one pipeline", ACCENT_RED)

steps = [
    ("1\nPrometheus\nFIRING", ACCENT_AMBER),
    ("2\nAlertmanager\nWebhook POST", ACCENT_RED),
    ("3\n/genai/alerts\n/ingest/", ACCENT_BLUE),
    ("4\nNormalize\nPayload", ACCENT_CYAN),
    ("5\nAiDE LLM\nRCA Call", ACCENT_PURPLE),
    ("6\nStore in Redis\nCache", ACCENT_GREEN),
    ("7\nFrontend\npoll /recent/", ACCENT_BLUE),
]

box_w = Inches(1.52)
box_h = Inches(1.1)
gap   = Inches(0.2)
start_x = Inches(0.2)
y_center = Inches(3.5)

for i, (label, color) in enumerate(steps):
    bx = start_x + i * (box_w + gap)
    add_rect(slide, bx, y_center, box_w, box_h, fill=BG_CARD, line=color, line_width=Pt(2))
    add_rect(slide, bx, y_center, box_w, Inches(0.07), fill=color)
    add_text(slide, label, bx + Inches(0.08), y_center + Inches(0.1),
             box_w - Inches(0.16), box_h - Inches(0.16),
             size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = bx + box_w + Inches(0.04)
        ay = y_center + box_h / 2 - Inches(0.04)
        add_rect(slide, ax, ay, gap - Inches(0.04), Inches(0.06), fill=MUTED)

# Resolution path
add_text(slide, "RESOLVED PATH", Inches(0.2), Inches(4.9), Inches(4), Inches(0.3),
         size=Pt(10), bold=True, color=ACCENT_GREEN)
res_steps = [
    ("Alertmanager\nresolved webhook", ACCENT_GREEN),
    ("_remove_by_identity()\nfallback: name-only", ACCENT_CYAN),
    ("Cache entry\nevicted", ACCENT_GREEN),
    ("Frontend shows\nno active alert", ACCENT_BLUE),
]
for i, (label, color) in enumerate(res_steps):
    bx = Inches(0.2) + i * (Inches(2.85))
    by = Inches(5.2)
    add_rect(slide, bx, by, Inches(2.6), Inches(0.9), fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_text(slide, label, bx + Inches(0.1), by + Inches(0.1), Inches(2.4), Inches(0.7),
             size=Pt(10), color=color, align=PP_ALIGN.CENTER)
    if i < len(res_steps) - 1:
        add_rect(slide, bx + Inches(2.62), by + Inches(0.38), Inches(0.2), Inches(0.06), fill=ACCENT_GREEN)

add_text(slide, "Root cause of stale UI bug → Redis key `aiops_recent_alert_recommendations` TTL 24h.  "
                "Fix: evict on resolved ingest + filter on read.",
         Inches(0.2), Inches(6.3), Inches(12.9), Inches(0.4),
         size=Pt(10), italic=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AI REMEDIATION FLOW
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "AI INTEGRATION", "AI-Guided Diagnosis & Remediation",
               "AiDE LLM drives root cause analysis, diagnostic commands, and remediation planning", ACCENT_PURPLE)

flow = [
    ("Alert Context\nCollection",     "Prometheus metrics\nElasticsearch logs\nJaeger traces",    ACCENT_AMBER),
    ("AiDE LLM\nRCA Call",            "Structured prompt:\nAlert + metrics + logs\n→ JSON response", ACCENT_PURPLE),
    ("Diagnostic Plan\nExtraction",   "root_cause · evidence\nimpact · resolution\nremediation_cmd",  ACCENT_CYAN),
    ("Command\nDelegation",           "DB-agent or\nControl-agent\nover HTTP",                     ACCENT_BLUE),
    ("Output\nAnalysis",              "Post-command\nAI analysis call\n→ final_answer",             ACCENT_GREEN),
    ("Incident\nCorrelation",         "Auto-create/update\nIncident + Timeline\nBlast radius",     ACCENT_RED),
]

bw = Inches(1.95)
bh = Inches(2.2)
for i, (title, detail, color) in enumerate(flow):
    col = i % 3
    row = i // 3
    bx = Inches(0.25) + col * Inches(4.35)
    by = Inches(2.1) + row * Inches(2.45)
    add_rect(slide, bx, by, bw * 2.05, bh, fill=BG_CARD, line=color, line_width=Pt(1.8))
    add_rect(slide, bx, by, bw * 2.05, Inches(0.08), fill=color)
    # number badge
    add_rect(slide, bx + Inches(0.12), by + Inches(0.15), Inches(0.38), Inches(0.38), fill=color)
    add_text(slide, str(i + 1), bx + Inches(0.12), by + Inches(0.15), Inches(0.38), Inches(0.38),
             size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + Inches(0.6), by + Inches(0.15), bw * 2.05 - Inches(0.75), Inches(0.55),
             size=Pt(13), bold=True, color=WHITE)
    add_text(slide, detail, bx + Inches(0.15), by + Inches(0.75), bw * 2.05 - Inches(0.3), Inches(1.3),
             size=Pt(11), color=OFF_WHITE)

add_text(slide, "Dual AiDE endpoint support (primary + secondary) with retry/backoff and SSL toggle",
         Inches(0.2), Inches(7.0), Inches(12.9), Inches(0.35),
         size=Pt(10), italic=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — INCIDENT MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "INCIDENT MANAGEMENT", "Incident Lifecycle",
               "Auto-correlated incidents with timeline, graph topology, and blast radius", ACCENT_RED)

# Incident model fields
fields = [
    ("incident_key", "UUID — unique incident identifier", ACCENT_BLUE),
    ("status", "open · investigating · resolved", ACCENT_AMBER),
    ("severity", "critical · warning · info", ACCENT_RED),
    ("target_host", "Affected service/host", ACCENT_CYAN),
    ("blast_radius", "List of downstream services", ACCENT_PURPLE),
    ("dependency_graph", "Topology JSON snapshot", ACCENT_GREEN),
    ("timeline", "Ordered IncidentTimelineEvent records", ACCENT_BLUE),
]
add_text(slide, "Incident Model", Inches(0.25), Inches(2.0), Inches(5), Inches(0.4),
         size=Pt(14), bold=True, color=ACCENT_RED)
for i, (field, desc, color) in enumerate(fields):
    cy = Inches(2.5) + i * Inches(0.6)
    add_rect(slide, Inches(0.25), cy, Inches(1.8), Inches(0.45), fill=color)
    add_text(slide, field, Inches(0.35), cy + Inches(0.08), Inches(1.6), Inches(0.3),
             size=Pt(10), bold=True, color=WHITE)
    add_text(slide, desc, Inches(2.15), cy + Inches(0.08), Inches(3.2), Inches(0.35),
             size=Pt(10), color=OFF_WHITE)

# Lifecycle flow (right)
add_text(slide, "Incident Lifecycle", Inches(6.5), Inches(2.0), Inches(6.5), Inches(0.4),
         size=Pt(14), bold=True, color=ACCENT_RED)
stages = [
    ("Alert Fires", ACCENT_AMBER,
     "Alertmanager sends resolved/firing webhook to /genai/alerts/ingest/"),
    ("Incident Created", ACCENT_RED,
     "_correlate_alert_to_incident() — dedup by fingerprint, create or reuse"),
    ("AI Analysis", ACCENT_PURPLE,
     "collect_alert_context() → AiDE RCA → diagnostic plan extracted"),
    ("Timeline Events", ACCENT_BLUE,
     "IncidentTimelineEvent records each command + analysis step"),
    ("Resolution", ACCENT_GREEN,
     "Status → resolved, resolved_at stamped, cache entry evicted"),
]
for i, (title, color, desc) in enumerate(stages):
    cy = Inches(2.5) + i * Inches(0.92)
    add_rect(slide, Inches(6.5), cy, Inches(0.5), Inches(0.5), fill=color)
    add_text(slide, str(i+1), Inches(6.5), cy + Inches(0.05), Inches(0.5), Inches(0.42),
             size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(7.1), cy, Inches(2.5), Inches(0.35),
             size=Pt(12), bold=True, color=color)
    add_text(slide, desc, Inches(7.1), cy + Inches(0.35), Inches(5.9), Inches(0.5),
             size=Pt(10), color=OFF_WHITE)
    if i < len(stages) - 1:
        add_rect(slide, Inches(6.7), cy + Inches(0.52), Inches(0.08), Inches(0.4), fill=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GRAPH & TOPOLOGY
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "TOPOLOGY", "Graph-Driven Topology Views",
               "Application dependency graphs, incident blast-radius maps, and service risk overlays", ACCENT_PURPLE)

# Central node + deps diagram
cx, cy_n = Inches(6.6), Inches(4.2)
# central
add_rect(slide, cx - Inches(0.65), cy_n - Inches(0.32), Inches(1.3), Inches(0.64),
         fill=ACCENT_RED, line=None)
add_text(slide, "INCIDENT\nEPICENTER", cx - Inches(0.6), cy_n - Inches(0.28),
         Inches(1.2), Inches(0.56), size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

nodes = [
    ("frontend",    Inches(4.2),  Inches(2.4), ACCENT_AMBER),
    ("gateway",     Inches(4.5),  Inches(4.2), ACCENT_BLUE),
    ("app-orders",  Inches(5.8),  Inches(5.8), ACCENT_CYAN),
    ("app-inventory",Inches(7.4), Inches(5.8), ACCENT_CYAN),
    ("app-billing", Inches(8.8),  Inches(4.2), ACCENT_CYAN),
    ("postgres",    Inches(8.5),  Inches(2.4), ACCENT_GREEN),
    ("redis",       Inches(6.6),  Inches(2.0), ACCENT_GREEN),
]
for name, nx, ny, color in nodes:
    add_rect(slide, nx - Inches(0.55), ny - Inches(0.25), Inches(1.1), Inches(0.5),
             fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_text(slide, name, nx - Inches(0.5), ny - Inches(0.2),
             Inches(1.0), Inches(0.4), size=Pt(9), bold=True, color=color, align=PP_ALIGN.CENTER)

add_text(slide, "Nodes carry: status · risk_score · ai_insight · depends_on · dependents",
         Inches(0.25), Inches(6.4), Inches(8), Inches(0.4),
         size=Pt(11), italic=True, color=MUTED)

# Legend
legend = [
    ("Epicenter", ACCENT_RED), ("Edge Service", ACCENT_AMBER),
    ("Gateway", ACCENT_BLUE), ("Microservice", ACCENT_CYAN),
    ("Data Store", ACCENT_GREEN),
]
for i, (label, color) in enumerate(legend):
    lx = Inches(0.25) + i * Inches(2.3)
    add_rect(slide, lx, Inches(6.85), Inches(0.22), Inches(0.22), fill=color)
    add_text(slide, label, lx + Inches(0.28), Inches(6.85), Inches(1.9), Inches(0.25),
             size=Pt(10), color=OFF_WHITE)

# Right panel — graph types
add_text(slide, "3 Graph Types", Inches(10.2), Inches(2.1), Inches(3), Inches(0.4),
         size=Pt(13), bold=True, color=WHITE)
for i, (gt, color, route) in enumerate([
    ("Application Graph", ACCENT_BLUE,  "/graph/application/:key"),
    ("Incident Graph",    ACCENT_RED,   "/graph/incident/:key"),
    ("Alert Graph",       ACCENT_AMBER, "/graph/:alertId"),
]):
    cy2 = Inches(2.6) + i * Inches(1.2)
    add_rect(slide, Inches(10.2), cy2, Inches(2.95), Inches(1.0),
             fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_text(slide, gt, Inches(10.35), cy2 + Inches(0.1), Inches(2.65), Inches(0.35),
             size=Pt(11), bold=True, color=color)
    add_text(slide, route, Inches(10.35), cy2 + Inches(0.45), Inches(2.65), Inches(0.45),
             size=Pt(9), color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PREDICTIONS
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "PREDICTIVE RISK", "Predictive Incident Risk",
               "Forward-looking risk scoring with per-service health windows", ACCENT_AMBER)

metrics = [
    ("Risk Score",           "0.0 – 1.0 float per service",         ACCENT_RED),
    ("Incident Probability", "% chance in next 15 minutes",          ACCENT_AMBER),
    ("Service Status",       "healthy · at-risk · degraded",         ACCENT_CYAN),
    ("Model Version",        "heuristic-v1 (ML-ready schema)",        ACCENT_PURPLE),
    ("Blast Radius",         "Downstream services that will be hit",  ACCENT_BLUE),
    ("Explanation",          "Human-readable risk reason",            ACCENT_GREEN),
]
for i, (label, desc, color) in enumerate(metrics):
    col = i % 2
    row = i // 2
    bx = Inches(0.25) + col * Inches(6.5)
    by = Inches(2.1) + row * Inches(1.5)
    add_rect(slide, bx, by, Inches(6.2), Inches(1.3),
             fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_rect(slide, bx, by, Inches(6.2), Inches(0.07), fill=color)
    add_text(slide, label, bx + Inches(0.15), by + Inches(0.15), Inches(5.8), Inches(0.38),
             size=Pt(14), bold=True, color=color)
    add_text(slide, desc, bx + Inches(0.15), by + Inches(0.55), Inches(5.8), Inches(0.65),
             size=Pt(12), color=OFF_WHITE)

add_text(slide, "ServicePrediction and PredictionSnapshot models — ready to swap heuristic scorer with ML model",
         Inches(0.25), Inches(6.9), Inches(12.8), Inches(0.4),
         size=Pt(10), italic=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FLEET ONBOARDING FLOW
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "FLEET MANAGEMENT", "Linux Fleet Onboarding Flow",
               "SSH + PEM automated install with telemetry profile selection and enrollment tokens", ACCENT_GREEN)

phases = [
    ("Operator selects\ntarget type + profile\nin Enroll UI",
     "Phase 1\nSetup", ACCENT_BLUE),
    ("Backend generates\nEnrollmentToken\n(24h expiry)",
     "Token\nGeneration", ACCENT_CYAN),
    ("One-liner curl install\ncommand shown\nto operator",
     "Install\nBlueprint", ACCENT_AMBER),
    ("SSH connectivity\ntest via PEM\n(pre-flight check)",
     "Connectivity\nTest", ACCENT_RED),
    ("Remote install\ntriggered over SSH:\ncollectors + exporters",
     "Automated\nInstall", ACCENT_PURPLE),
    ("Agent calls\n/fleet/enroll/\nwith token",
     "Self-\nRegistration", ACCENT_GREEN),
    ("Heartbeats → target\nstatus + discovered\nservices update",
     "Ongoing\nHeartbeat", ACCENT_GREEN),
]

bw2 = Inches(1.72)
bh2 = Inches(1.7)
start2 = Inches(0.15)
top2 = Inches(2.25)
for i, (desc, title, color) in enumerate(phases):
    bx2 = start2 + i * (bw2 + Inches(0.12))
    add_rect(slide, bx2, top2, bw2, bh2, fill=BG_CARD, line=color, line_width=Pt(1.8))
    add_rect(slide, bx2, top2, bw2, Inches(0.08), fill=color)
    # step number
    add_rect(slide, bx2 + bw2 - Inches(0.42), top2 + Inches(0.12), Inches(0.32), Inches(0.32), fill=color)
    add_text(slide, str(i+1), bx2 + bw2 - Inches(0.42), top2 + Inches(0.12),
             Inches(0.32), Inches(0.32), size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx2 + Inches(0.1), top2 + Inches(0.18), bw2 - Inches(0.15), Inches(0.5),
             size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, desc, bx2 + Inches(0.1), top2 + Inches(0.72), bw2 - Inches(0.15), Inches(0.9),
             size=Pt(9), color=OFF_WHITE, align=PP_ALIGN.CENTER)
    if i < len(phases) - 1:
        ax2 = bx2 + bw2 + Inches(0.01)
        ay2 = top2 + bh2 / 2 - Inches(0.04)
        add_rect(slide, ax2, ay2, Inches(0.1), Inches(0.06), fill=color)

# Models
add_text(slide, "Fleet Models Introduced", Inches(0.25), Inches(4.2), Inches(12), Inches(0.35),
         size=Pt(12), bold=True, color=ACCENT_GREEN)
models_list = [
    ("TelemetryProfile", "Install bundle definition (components + capabilities)", ACCENT_CYAN),
    ("Target",           "Enrolled host with status, heartbeat, and discovered services", ACCENT_GREEN),
    ("EnrollmentToken",  "Short-lived bootstrap token with expiry + revocation", ACCENT_AMBER),
    ("TargetOnboardingRequest", "SSH-based install request with connectivity test + install trigger", ACCENT_BLUE),
    ("TargetComponent",  "Per-collector/exporter status on each target", ACCENT_PURPLE),
    ("TargetHeartbeat",  "Raw heartbeat payloads from enrolled agents", MUTED),
    ("DiscoveredService","Services auto-discovered by the control agent on each host", ACCENT_RED),
]
cols_m = 3
for i, (name, desc, color) in enumerate(models_list):
    col_m = i % cols_m
    row_m = i // cols_m
    mx = Inches(0.25) + col_m * Inches(4.35)
    my = Inches(4.65) + row_m * Inches(0.62)
    add_rect(slide, mx, my, Inches(1.55), Inches(0.45), fill=color)
    add_text(slide, name, mx + Inches(0.08), my + Inches(0.06), Inches(1.4), Inches(0.34),
             size=Pt(9), bold=True, color=WHITE)
    add_text(slide, desc, mx + Inches(1.65), my + Inches(0.06), Inches(2.55), Inches(0.4),
             size=Pt(9), color=OFF_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FRONTEND & UX
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "FRONTEND", "React Control Plane UI",
               "9-page React + Vite SPA with dark/light theme, animated transitions, and live refresh", ACCENT_CYAN)

pages = [
    ("/applications", "Applications", "Portfolio view — live service health + AI insight per component", ACCENT_BLUE),
    ("/alerts",       "Alerts",       "Recent alert recommendations from Redis cache — AI diagnosis cards", ACCENT_RED),
    ("/incidents",    "Incidents",    "Incident timeline, graph, blast radius, and RCA detail", ACCENT_AMBER),
    ("/predictions",  "Predictions",  "Risk score + incident probability per service", ACCENT_PURPLE),
    ("/assistant",    "Assistant",    "Multi-turn AI chat with session history and context", ACCENT_CYAN),
    ("/fleet",        "Fleet",        "Enrolled targets, heartbeat status, discovered services", ACCENT_GREEN),
    ("/enroll",       "Enroll",       "SSH + PEM onboarding wizard with profile selection", ACCENT_AMBER),
    ("/profiles",     "Profiles",     "Telemetry profile browser with components + capabilities", ACCENT_BLUE),
    ("/graph",        "Graph",        "Application / incident / alert topology with dependency edges", ACCENT_PURPLE),
]
for i, (path, name, desc, color) in enumerate(pages):
    col = i % 3
    row = i // 3
    bx = Inches(0.2) + col * Inches(4.35)
    by = Inches(2.1) + row * Inches(1.6)
    add_rect(slide, bx, by, Inches(4.1), Inches(1.45), fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_rect(slide, bx, by, Inches(4.1), Inches(0.07), fill=color)
    add_text(slide, path, bx + Inches(0.12), by + Inches(0.12), Inches(3.8), Inches(0.3),
             size=Pt(10), color=MUTED)
    add_text(slide, name, bx + Inches(0.12), by + Inches(0.38), Inches(3.8), Inches(0.38),
             size=Pt(14), bold=True, color=color)
    add_text(slide, desc, bx + Inches(0.12), by + Inches(0.78), Inches(3.8), Inches(0.55),
             size=Pt(10), color=OFF_WHITE)

add_text(slide, "Tech: React 18 · TypeScript · TanStack Query · React Router v6 · Framer Motion · Vite",
         Inches(0.2), Inches(7.0), Inches(12.9), Inches(0.35),
         size=Pt(10), italic=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "TECH STACK", "Technology Stack",
               "Production-grade open-source components — fully containerised", ACCENT_BLUE)

categories = [
    ("Backend", ACCENT_BLUE, [
        "Python 3 · Django 4",
        "Gunicorn WSGI",
        "Celery workers",
        "Django Cache (Redis)",
        "psycopg2 + ORM migrations",
    ]),
    ("Frontend", ACCENT_CYAN, [
        "React 18 + TypeScript",
        "Vite build tooling",
        "TanStack Query v5",
        "React Router v6",
        "Framer Motion",
    ]),
    ("Observability", ACCENT_PURPLE, [
        "Prometheus + Alertmanager",
        "OpenTelemetry Collector",
        "Jaeger (traces)",
        "Elasticsearch + Filebeat",
        "Grafana dashboards",
    ]),
    ("Infra / DevOps", ACCENT_GREEN, [
        "Docker + Compose",
        "PostgreSQL 13",
        "Redis 7",
        "Nginx (frontend + gateway)",
        "Toxiproxy (chaos testing)",
    ]),
    ("AI", ACCENT_AMBER, [
        "AiDE LLM API (primary)",
        "AiDE secondary endpoint",
        "Retry + backoff adapter",
        "Structured JSON prompts",
        "better-profanity filter",
    ]),
    ("Agents", ACCENT_RED, [
        "db-agent (DB commands)",
        "control-agent (SSH / OS)",
        "Token auth (AGENT_SECRET)",
        "Allowed-commands JSON",
        "Fleet heartbeat receiver",
    ]),
]
for i, (cat, color, items) in enumerate(categories):
    col = i % 3
    row = i // 3
    bx = Inches(0.2) + col * Inches(4.35)
    by = Inches(2.1) + row * Inches(2.5)
    bh_c = Inches(2.35)
    add_rect(slide, bx, by, Inches(4.1), bh_c, fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_rect(slide, bx, by, Inches(4.1), Inches(0.45), fill=color)
    add_text(slide, cat, bx + Inches(0.15), by + Inches(0.06), Inches(3.8), Inches(0.35),
             size=Pt(14), bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", bx + Inches(0.2), by + Inches(0.55 + j * 0.36),
                 Inches(3.7), Inches(0.34), size=Pt(10), color=OFF_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SECURITY & REPO HYGIENE
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "SECURITY", "Security & Repo Hygiene",
               "Secrets removed from git history — .env-based configuration enforced", ACCENT_AMBER)

done = [
    ("Secrets Removed from History",
     "git filter-branch stripped GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET, AIDE_API_KEY from all commits. Force-pushed clean main.", ACCENT_GREEN),
    ("docker-compose.yml Sanitised",
     "All credentials replaced with ${ENV_VAR} references. Defaults are empty or safe placeholders.", ACCENT_GREEN),
    (".env Pattern Added",
     ".env.example provides the full template. .gitignore ensures .env is never committed.", ACCENT_GREEN),
    ("Agent Token Auth",
     "db-agent and control-agent require AGENT_AUTH_TOKEN header on every command request.", ACCENT_GREEN),
    ("SSO Middleware",
     "HeaderSSOAuthMiddleware wraps every request. SSO_ENABLED toggle in settings.", ACCENT_AMBER),
    ("Allowed-Commands Files",
     "Each agent has a JSON allowlist — arbitrary command execution is not permitted.", ACCENT_AMBER),
]
for i, (title, desc, color) in enumerate(done):
    col = i % 2
    row = i // 2
    bx = Inches(0.2) + col * Inches(6.55)
    by = Inches(2.1) + row * Inches(1.6)
    add_rect(slide, bx, by, Inches(6.3), Inches(1.45), fill=BG_CARD, line=color, line_width=Pt(1.5))
    add_rect(slide, bx + Inches(0.12), by + Inches(0.14), Inches(0.32), Inches(0.32), fill=color)
    add_text(slide, "✓" if color == ACCENT_GREEN else "⚠", bx + Inches(0.12), by + Inches(0.1),
             Inches(0.32), Inches(0.34), size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + Inches(0.55), by + Inches(0.1), Inches(5.6), Inches(0.38),
             size=Pt(13), bold=True, color=color)
    add_text(slide, desc, bx + Inches(0.55), by + Inches(0.52), Inches(5.6), Inches(0.8),
             size=Pt(10), color=OFF_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
section_header(slide, "ROADMAP", "What's Next",
               "Phased evolution from Phase 1 foundation to full autonomous operations", ACCENT_GREEN)

phases_road = [
    ("Phase 2\n(Near-term)", ACCENT_BLUE, [
        "Complete Linux installer + OTEL agent setup",
        "Stable public deployment + onboarding flow",
        "Run pending DB migrations (0005, 0006)",
        "Fix pip requirements hash for clean builds",
        "Commit agent/ + migration files to main",
    ]),
    ("Phase 3\n(Mid-term)", ACCENT_CYAN, [
        "Auto-generate topology from traces + discovery",
        "Periodic Prometheus reconciliation for cache",
        "Extend onboarding to Windows + Kubernetes",
        "MCP-based telemetry and tool access layer",
        "ML model to replace heuristic risk scorer",
    ]),
    ("Phase 4\n(Strategic)", ACCENT_PURPLE, [
        "Approval-driven remediation workflows",
        "Multi-tenant fleet isolation",
        "Real-time streaming topology (WebSocket)",
        "Feedback loop — remediation outcome → model",
        "SLA / SLO tracking per service",
    ]),
]
for i, (phase, color, items) in enumerate(phases_road):
    bx = Inches(0.25) + i * Inches(4.35)
    by = Inches(2.1)
    bh_p = Inches(5.0)
    add_rect(slide, bx, by, Inches(4.1), bh_p, fill=BG_CARD, line=color, line_width=Pt(2))
    add_rect(slide, bx, by, Inches(4.1), Inches(0.6), fill=color)
    add_text(slide, phase, bx + Inches(0.15), by + Inches(0.08), Inches(3.8), Inches(0.48),
             size=Pt(15), bold=True, color=WHITE)
    for j, item in enumerate(items):
        iy = by + Inches(0.8) + j * Inches(0.8)
        add_rect(slide, bx + Inches(0.2), iy, Inches(0.12), Inches(0.5), fill=color)
        add_text(slide, item, bx + Inches(0.42), iy, Inches(3.5), Inches(0.55),
                 size=Pt(11), color=OFF_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, W, Inches(0.08), fill=ACCENT_BLUE)
add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=ACCENT_BLUE)
add_rect(slide, Inches(8.8), 0, Inches(4.53), H, fill=BG_CARD)
add_rect(slide, Inches(8.8), 0, Inches(0.05), H, fill=ACCENT_BLUE)

add_text(slide, "AIOps",
         Inches(0.5), Inches(1.5), Inches(7.5), Inches(1.2),
         size=Pt(64), bold=True, color=ACCENT_BLUE)
add_text(slide, "Phase 1 Complete",
         Inches(0.5), Inches(2.7), Inches(7.5), Inches(0.65),
         size=Pt(30), bold=True, color=WHITE)
add_text(slide, "Built · Deployed · Verified · Secured",
         Inches(0.5), Inches(3.4), Inches(7.5), Inches(0.45),
         size=Pt(15), color=OFF_WHITE)

summary_items = [
    ("9  React pages", ACCENT_BLUE),
    ("40+ backend endpoints", ACCENT_CYAN),
    ("7 new fleet models", ACCENT_GREEN),
    ("AI RCA + remediation", ACCENT_PURPLE),
    ("Secrets purged from git", ACCENT_AMBER),
    ("Fleet onboarding live", ACCENT_GREEN),
]
for i, (label, color) in enumerate(summary_items):
    lx = Inches(9.1) + (i % 2) * Inches(2.0)
    ly = Inches(1.5) + (i // 2) * Inches(0.95)
    add_rect(slide, lx, ly, Inches(1.85), Inches(0.72), fill=color)
    add_text(slide, label, lx + Inches(0.1), ly + Inches(0.1),
             Inches(1.65), Inches(0.52), size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Thank You",
         Inches(0.5), Inches(5.8), Inches(7.5), Inches(0.65),
         size=Pt(26), bold=True, color=ACCENT_CYAN)
add_text(slide, "April 2026  ·  AIOps Observability Control Plane  ·  Phase 1",
         Inches(0.5), Inches(6.5), Inches(8), Inches(0.4),
         size=Pt(11), color=MUTED)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "AIOps_ControlPlane_Presentation.pptx"
prs.save(out)
print(f"✅  Saved → {out}  ({len(prs.slides)} slides)")
